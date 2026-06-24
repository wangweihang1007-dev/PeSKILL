from __future__ import annotations

import argparse
import json
from pathlib import Path


SUPPORTED = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def extract_pdf(path: Path) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return {
        "type": "pdf",
        "pages": [
            {"page": index, "text": page.extract_text() or ""}
            for index, page in enumerate(reader.pages, start=1)
        ],
    }


def extract_docx(path: Path) -> dict:
    from docx import Document

    document = Document(str(path))
    return {
        "type": "docx",
        "paragraphs": [p.text for p in document.paragraphs if p.text.strip()],
        "tables": [
            [[cell.text for cell in row.cells] for row in table.rows]
            for table in document.tables
        ],
    }


def extract_pptx(path: Path) -> dict:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text for cell in row.cells))
        slides.append({"slide": index, "text": "\n".join(texts)})
    return {"type": "pptx", "slides": slides}


def extract_text(path: Path) -> dict:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return {
                "type": path.suffix.lower().lstrip("."),
                "encoding": encoding,
                "text": path.read_text(encoding=encoding),
            }
        except UnicodeDecodeError:
            continue
    raise ValueError(f"Cannot decode text file: {path}")


def extract(path: Path) -> dict:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        content = extract_pdf(path)
    elif suffix == ".docx":
        content = extract_docx(path)
    elif suffix == ".pptx":
        content = extract_pptx(path)
    elif suffix in {".txt", ".md"}:
        content = extract_text(path)
    else:
        raise ValueError(f"Unsupported file type: {path}")
    return {"path": str(path.resolve()), **content}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input", type=Path)
    parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args()

    paths = (
        sorted(p for p in args.input.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED)
        if args.input.is_dir()
        else [args.input]
    )
    args.output_dir.mkdir(parents=True, exist_ok=True)
    manifest = []
    for path in paths:
        payload = extract(path)
        output = args.output_dir / f"{path.stem}.json"
        output.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        manifest.append({"source": str(path.resolve()), "output": str(output.resolve())})
    manifest_path = args.output_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"files": len(manifest), "manifest": str(manifest_path.resolve())}, ensure_ascii=False))


if __name__ == "__main__":
    main()
